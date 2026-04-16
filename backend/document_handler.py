"""Shared handlers for extracting and replacing text in supported documents."""
from __future__ import annotations

from dataclasses import dataclass
import html
import os
from pathlib import Path
import re
from typing import Any, Dict, List, Optional

from docx import Document
from docx.table import _Cell as DocxCell
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape
from pptx.table import _Cell as PptxCell
from pptx.text.text import _Paragraph
from pptx.util import Pt


@dataclass
class TextElement:
    element_id: str
    text: str


@dataclass(frozen=True)
class TranslationOptions:
    side_by_side: bool = False
    pdf_output_format: str = "pdf"
    spreadsheet_mode: str = "in_place"


def _get_run_font_attributes(runs) -> tuple[Optional[str], Optional[Any]]:
    """Return the first available font name and size from a collection of runs."""
    for run in runs:
        font = getattr(run, "font", None)
        if not font:
            continue
        name = getattr(font, "name", None)
        size = getattr(font, "size", None)
        if name or size:
            return name, size
    return None, None


class DocumentHandler:
    """Base protocol for document handlers."""

    def extract_text(self) -> List[TextElement]:  # pragma: no cover - protocol
        raise NotImplementedError

    def apply_translations(
        self,
        translations: Dict[str, str],
        output_path: Path,
        font_name: Optional[str] = None,
        options: Optional[TranslationOptions] = None,
        original_texts: Optional[Dict[str, str]] = None,
    ) -> Path:  # pragma: no cover - protocol
        raise NotImplementedError


class PowerPointHandler(DocumentHandler):
    """Extracts and applies translations for PowerPoint presentations."""

    BODY_FONT_NAME = "Arial"
    BODY_FONT_SIZE_PT = 24

    def __init__(self, presentation_path: Path):
        self.presentation_path = Path(presentation_path)

    def extract_text(self) -> List[TextElement]:
        presentation = Presentation(self.presentation_path)
        elements: List[TextElement] = []

        for slide_index, slide in enumerate(presentation.slides):
            for shape_index, shape in enumerate(slide.shapes):
                elements.extend(self._extract_from_shape(slide_index, shape_index, shape))

        return elements

    def apply_translations(
        self,
        translations: Dict[str, str],
        output_path: Path,
        font_name: Optional[str] = None,
        options: Optional[TranslationOptions] = None,
        original_texts: Optional[Dict[str, str]] = None,
    ) -> Path:
        presentation = Presentation(self.presentation_path)
        original_texts = original_texts or {}
        options = options or TranslationOptions()

        for slide_index, slide in enumerate(presentation.slides):
            for shape_index, shape in enumerate(slide.shapes):
                self._apply_to_shape(
                    slide_index,
                    shape_index,
                    shape,
                    translations,
                    options.side_by_side,
                    original_texts,
                )

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        return output_path

    # PowerPoint helpers --------------------------------------------------
    def _extract_from_shape(
        self, slide_index: int, shape_index: int, shape: BaseShape
    ) -> List[TextElement]:
        elements: List[TextElement] = []

        if getattr(shape, "has_text_frame", False):
            text_frame = shape.text_frame
            for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
                text = paragraph.text.strip()
                if text:
                    element_id = self._paragraph_id(slide_index, shape_index, paragraph_index)
                    elements.append(TextElement(element_id=element_id, text=text))

        if getattr(shape, "has_table", False):
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    text = cell.text.strip()
                    if text:
                        element_id = self._cell_id(slide_index, shape_index, row_idx, col_idx)
                        elements.append(TextElement(element_id=element_id, text=text))

        return elements

    def _apply_to_shape(
        self,
        slide_index: int,
        shape_index: int,
        shape: BaseShape,
        translations: Dict[str, str],
        side_by_side: bool = False,
        original_texts: Optional[Dict[str, str]] = None,
    ) -> None:
        original_texts = original_texts or {}
        if getattr(shape, "has_text_frame", False):
            text_frame = shape.text_frame
            for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
                element_id = self._paragraph_id(slide_index, shape_index, paragraph_index)
                translated = translations.get(element_id)
                if translated is not None:
                    original_text = original_texts.get(element_id, paragraph.text.strip())
                    if side_by_side:
                        formatted_text = self._format_side_by_side(original_text, translated)
                        self._replace_ppt_paragraph(paragraph, formatted_text, font_size_pt=16)
                    else:
                        self._replace_ppt_paragraph(paragraph, translated)

        if getattr(shape, "has_table", False):
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    element_id = self._cell_id(slide_index, shape_index, row_idx, col_idx)
                    translated = translations.get(element_id)
                    if translated is not None:
                        original_text = original_texts.get(element_id, cell.text.strip())
                        if side_by_side:
                            formatted_text = self._format_side_by_side(original_text, translated)
                            self._replace_ppt_cell(cell, formatted_text, font_size_pt=16)
                        else:
                            self._replace_ppt_cell(cell, translated)

    @classmethod
    def _replace_ppt_paragraph(cls, paragraph: _Paragraph, new_text: str, font_size_pt: Optional[int] = None) -> None:
        paragraph.clear()
        run = paragraph.add_run()
        run.text = new_text
        cls._apply_body_text_format(paragraph, run, font_size_pt=font_size_pt)

    @classmethod
    def _replace_ppt_cell(cls, cell: PptxCell, new_text: str, font_size_pt: Optional[int] = None) -> None:
        text_frame = cell.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.clear()
        run = paragraph.add_run()
        run.text = new_text
        cls._apply_body_text_format(paragraph, run, font_size_pt=font_size_pt)

    @staticmethod
    def _paragraph_id(slide_index: int, shape_index: int, paragraph_index: int) -> str:
        return f"ppt_s{slide_index}_sh{shape_index}_p{paragraph_index}"

    @staticmethod
    def _cell_id(slide_index: int, shape_index: int, row_index: int, col_index: int) -> str:
        return f"ppt_s{slide_index}_sh{shape_index}_c{row_index}_{col_index}"

    @staticmethod
    def _format_side_by_side(original: str, translated: str) -> str:
        """Format text as 'original | translated' for side-by-side display."""
        return f"{original} | {translated}"

    @classmethod
    def _apply_body_text_format(cls, paragraph: _Paragraph, run, font_size_pt: Optional[int] = None) -> None:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1
        run.font.name = cls.BODY_FONT_NAME
        run.font.size = Pt(font_size_pt if font_size_pt is not None else cls.BODY_FONT_SIZE_PT)


class WordHandler(DocumentHandler):
    """Handles extraction and replacement for Word documents."""

    def __init__(self, document_path: Path):
        self.document_path = Path(document_path)

    def extract_text(self) -> List[TextElement]:
        document = Document(self.document_path)
        elements: List[TextElement] = []

        for paragraph_index, paragraph in enumerate(document.paragraphs):
            text = paragraph.text.strip()
            if text:
                elements.append(TextElement(self._paragraph_id(paragraph_index), text))

        for table_index, table in enumerate(document.tables):
            for row_index, row in enumerate(table.rows):
                for col_index, cell in enumerate(row.cells):
                    text = cell.text.strip()
                    if text:
                        elements.append(
                            TextElement(self._cell_id(table_index, row_index, col_index), text)
                        )

        return elements

    def apply_translations(
        self,
        translations: Dict[str, str],
        output_path: Path,
        font_name: Optional[str] = None,
        options: Optional[TranslationOptions] = None,
        original_texts: Optional[Dict[str, str]] = None,
    ) -> Path:
        document = Document(self.document_path)

        for paragraph_index, paragraph in enumerate(document.paragraphs):
            element_id = self._paragraph_id(paragraph_index)
            translated = translations.get(element_id)
            if translated is not None:
                self._replace_paragraph(paragraph, translated, font_name)

        for table_index, table in enumerate(document.tables):
            for row_index, row in enumerate(table.rows):
                for col_index, cell in enumerate(row.cells):
                    element_id = self._cell_id(table_index, row_index, col_index)
                    translated = translations.get(element_id)
                    if translated is not None:
                        self._replace_cell(cell, translated, font_name)

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        document.save(output_path)
        return output_path

    @staticmethod
    def _replace_paragraph(paragraph: Paragraph, new_text: str, font_name: Optional[str]) -> None:
        original_name, original_size = _get_run_font_attributes(paragraph.runs)
        if not original_size and getattr(paragraph, "style", None) and getattr(paragraph.style.font, "size", None):
            original_size = paragraph.style.font.size
        if not original_name and getattr(paragraph, "style", None) and getattr(paragraph.style.font, "name", None):
            original_name = paragraph.style.font.name
        for run in list(paragraph.runs):
            paragraph._p.remove(run._r)
        run = paragraph.add_run(new_text)
        if font_name:
            run.font.name = font_name
        elif original_name:
            run.font.name = original_name
        if original_size:
            run.font.size = original_size

    @staticmethod
    def _replace_cell(cell: DocxCell, new_text: str, font_name: Optional[str]) -> None:
        original_name: Optional[str] = None
        original_size: Optional[Any] = None
        for paragraph in cell.paragraphs:
            name, size = _get_run_font_attributes(paragraph.runs)
            if not original_name and name:
                original_name = name
            if not original_size and size:
                original_size = size
            if original_name and original_size:
                break
        if not original_size and cell.paragraphs:
            style_font = getattr(cell.paragraphs[0].style, "font", None)
            if style_font and getattr(style_font, "size", None):
                original_size = style_font.size
            if not original_name and style_font and getattr(style_font, "name", None):
                original_name = style_font.name
        cell.text = ""
        paragraph = cell.paragraphs[0] if cell.paragraphs else cell.add_paragraph("")
        for run in list(paragraph.runs):
            paragraph._p.remove(run._r)
        run = paragraph.add_run(new_text)
        if font_name:
            run.font.name = font_name
        elif original_name:
            run.font.name = original_name
        if original_size:
            run.font.size = original_size

    @staticmethod
    def _paragraph_id(paragraph_index: int) -> str:
        return f"doc_p{paragraph_index}"

    @staticmethod
    def _cell_id(table_index: int, row_index: int, col_index: int) -> str:
        return f"doc_t{table_index}_r{row_index}_c{col_index}"


class ExcelHandler(DocumentHandler):
    """Handles extraction and replacement for Excel workbooks."""

    def __init__(self, workbook_path: Path):
        self.workbook_path = Path(workbook_path)

    def extract_text(self) -> List[TextElement]:
        workbook = load_workbook(self.workbook_path, data_only=True)
        elements: List[TextElement] = []

        for sheet_index, sheet in enumerate(workbook.worksheets):
            for row in sheet.iter_rows():
                for cell in row:
                    value = cell.value
                    if isinstance(value, str):
                        text = value.strip()
                        if text:
                            elements.append(
                                TextElement(self._cell_id(sheet_index, cell.row, cell.column), text)
                            )
        return elements

    def apply_translations(
        self,
        translations: Dict[str, str],
        output_path: Path,
        font_name: Optional[str] = None,
        options: Optional[TranslationOptions] = None,
        original_texts: Optional[Dict[str, str]] = None,
    ) -> Path:
        options = options or TranslationOptions()
        workbook = load_workbook(self.workbook_path)
        original_sheets = list(workbook.worksheets)

        if options.spreadsheet_mode == "in_place":
            for sheet_index, sheet in enumerate(original_sheets):
                self._apply_to_sheet(sheet_index, sheet, translations, font_name=font_name)
        elif options.spreadsheet_mode == "new_sheet":
            for sheet_index, sheet in enumerate(original_sheets):
                translated_sheet = workbook.copy_worksheet(sheet)
                translated_sheet.title = self._translated_sheet_title(sheet.title, workbook.sheetnames)
                workbook._sheets.remove(translated_sheet)
                original_position = workbook._sheets.index(sheet)
                workbook._sheets.insert(original_position + 1, translated_sheet)
                self._apply_to_sheet(sheet_index, translated_sheet, translations, font_name=font_name)
        else:
            raise ValueError(f"Unsupported spreadsheet mode: {options.spreadsheet_mode}")

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        workbook.save(output_path)
        return output_path

    def _apply_to_sheet(
        self,
        sheet_index: int,
        sheet,
        translations: Dict[str, str],
        font_name: Optional[str] = None,
    ) -> None:
        for row in sheet.iter_rows():
            for cell in row:
                element_id = self._cell_id(sheet_index, cell.row, cell.column)
                translated = translations.get(element_id)
                if translated is not None:
                    cell.value = translated
                    if font_name:
                        cell.font = cell.font.copy(name=font_name)

    @staticmethod
    def _translated_sheet_title(base_title: str, existing_titles: List[str]) -> str:
        suffix = " (Translated)"
        max_base_length = 31 - len(suffix)
        candidate = f"{base_title[:max_base_length]}{suffix}"
        if candidate not in existing_titles:
            return candidate

        index = 2
        while True:
            numbered_suffix = f" (Translated {index})"
            max_base_length = 31 - len(numbered_suffix)
            candidate = f"{base_title[:max_base_length]}{numbered_suffix}"
            if candidate not in existing_titles:
                return candidate
            index += 1

    @staticmethod
    def _cell_id(sheet_index: int, row_index: int, column_index: int) -> str:
        return f"xls_s{sheet_index}_r{row_index}_c{column_index}"


class PdfHandler(DocumentHandler):
    """Extracts text from PDFs and renders translated output as PDF or DOCX."""

    def __init__(self, pdf_path: Path):
        self.pdf_path = Path(pdf_path)

    def extract_text(self) -> List[TextElement]:
        elements: List[TextElement] = []
        for page_index, blocks in enumerate(self._extract_page_blocks()):
            for block_index, text in enumerate(blocks):
                elements.append(TextElement(self._block_id(page_index, block_index), text))
        return elements

    def apply_translations(
        self,
        translations: Dict[str, str],
        output_path: Path,
        font_name: Optional[str] = None,
        options: Optional[TranslationOptions] = None,
        original_texts: Optional[Dict[str, str]] = None,
    ) -> Path:
        options = options or TranslationOptions()
        if options.pdf_output_format == "docx":
            return self._render_docx(translations, output_path, font_name=font_name)
        if options.pdf_output_format == "pdf":
            return self._render_pdf(translations, output_path, font_name=font_name)
        raise ValueError(f"Unsupported PDF output format: {options.pdf_output_format}")

    def _extract_page_blocks(self) -> List[List[str]]:
        reader = PdfReader(self.pdf_path)
        pages: List[List[str]] = []
        for page in reader.pages:
            raw_text = page.extract_text() or ""
            pages.append(self._split_pdf_text(raw_text))
        return pages

    def _render_docx(
        self,
        translations: Dict[str, str],
        output_path: Path,
        font_name: Optional[str] = None,
    ) -> Path:
        document = Document()
        page_blocks = self._extract_page_blocks()

        for page_index, blocks in enumerate(page_blocks):
            if len(page_blocks) > 1:
                heading = document.add_heading(f"Page {page_index + 1}", level=2)
                self._apply_font_to_docx_paragraph(heading, font_name)
            for block_index, original_text in enumerate(blocks):
                translated = translations.get(self._block_id(page_index, block_index), original_text)
                paragraph = document.add_paragraph(translated)
                self._apply_font_to_docx_paragraph(paragraph, font_name)
            if page_index < len(page_blocks) - 1:
                document.add_page_break()

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        document.save(output_path)
        return output_path

    def _render_pdf(
        self,
        translations: Dict[str, str],
        output_path: Path,
        font_name: Optional[str] = None,
    ) -> Path:
        try:
            from reportlab.lib.enums import TA_LEFT
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import mm
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.cidfonts import UnicodeCIDFont
            from reportlab.pdfbase.ttfonts import TTFont
            from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer
        except ModuleNotFoundError as exc:  # pragma: no cover - depends on runtime env
            raise RuntimeError(
                "PDF output requires the reportlab package. Run pip install -r backend/requirements.txt."
            ) from exc

        page_blocks = self._extract_page_blocks()
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        resolved_font_name = "Helvetica"
        font_path = self._resolve_pdf_font_path(font_name)
        if font_path:
            registered_name = "TranslatePdfFont"
            if registered_name not in pdfmetrics.getRegisteredFontNames():
                pdfmetrics.registerFont(TTFont(registered_name, str(font_path)))
            resolved_font_name = registered_name
        elif self._contains_cjk_text(translations.values()):
            cid_font_name = "STSong-Light"
            if cid_font_name not in pdfmetrics.getRegisteredFontNames():
                pdfmetrics.registerFont(UnicodeCIDFont(cid_font_name))
            resolved_font_name = cid_font_name

        styles = getSampleStyleSheet()
        body_style = ParagraphStyle(
            "TranslatedBody",
            parent=styles["BodyText"],
            fontName=resolved_font_name,
            fontSize=11,
            leading=16,
            alignment=TA_LEFT,
            spaceAfter=10,
            wordWrap="CJK",
        )
        heading_style = ParagraphStyle(
            "TranslatedHeading",
            parent=styles["Heading2"],
            fontName=resolved_font_name,
            fontSize=14,
            leading=18,
            spaceAfter=12,
            wordWrap="CJK",
        )

        story = []
        for page_index, blocks in enumerate(page_blocks):
            if len(page_blocks) > 1:
                story.append(Paragraph(html.escape(f"Page {page_index + 1}"), heading_style))
            for block_index, original_text in enumerate(blocks):
                translated = translations.get(self._block_id(page_index, block_index), original_text)
                story.append(Paragraph(self._to_pdf_markup(translated), body_style))
                story.append(Spacer(1, 3 * mm))
            if page_index < len(page_blocks) - 1:
                story.append(PageBreak())

        document = SimpleDocTemplate(
            str(output_path),
            pagesize=A4,
            leftMargin=18 * mm,
            rightMargin=18 * mm,
            topMargin=18 * mm,
            bottomMargin=18 * mm,
        )
        document.build(story)
        return output_path

    @staticmethod
    def _split_pdf_text(text: str) -> List[str]:
        normalized = text.replace("\r\n", "\n").replace("\r", "\n")
        blocks: List[str] = []
        for block in re.split(r"\n\s*\n+", normalized):
            lines = [line.strip() for line in block.split("\n") if line.strip()]
            if lines:
                blocks.append(" ".join(lines))
        return blocks

    @staticmethod
    def _block_id(page_index: int, block_index: int) -> str:
        return f"pdf_p{page_index}_b{block_index}"

    @staticmethod
    def _apply_font_to_docx_paragraph(paragraph: Paragraph, font_name: Optional[str]) -> None:
        if not font_name:
            return
        for run in paragraph.runs:
            run.font.name = font_name

    @staticmethod
    def _to_pdf_markup(text: str) -> str:
        return html.escape(text).replace("\n", "<br/>")

    @staticmethod
    def _resolve_pdf_font_path(font_name: Optional[str]) -> Optional[Path]:
        configured_path = os.getenv("PDF_FONT_PATH")
        if configured_path:
            candidate = Path(configured_path)
            if candidate.exists():
                return candidate

        candidate_names = [font_name] if font_name else []
        candidate_names.extend(
            [
                "NotoSansCJK-Regular.ttc",
                "NotoSansCJKsc-Regular.otf",
                "NotoSansSC-Regular.ttf",
                "SourceHanSansSC-Regular.otf",
                "PingFang.ttc",
                "Arial Unicode.ttf",
                "Arial Unicode MS.ttf",
                "DejaVuSans.ttf",
            ]
        )
        search_dirs = [
            Path("/System/Library/Fonts"),
            Path("/Library/Fonts"),
            Path("/usr/share/fonts"),
            Path("/usr/local/share/fonts"),
            Path.home() / ".fonts",
        ]

        for name in candidate_names:
            if not name:
                continue
            direct = Path(name)
            if direct.exists():
                return direct
            for directory in search_dirs:
                if not directory.exists():
                    continue
                for match in directory.rglob("*"):
                    if match.is_file() and match.name.lower() == Path(name).name.lower():
                        return match
        return None

    @staticmethod
    def _contains_cjk_text(values) -> bool:
        for value in values:
            if re.search(r"[\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff]", value or ""):
                return True
        return False


def get_document_handler(path: Path) -> DocumentHandler:
    suffix = Path(path).suffix.lower()
    if suffix == ".pptx":
        return PowerPointHandler(path)
    if suffix == ".docx":
        return WordHandler(path)
    if suffix == ".xlsx":
        return ExcelHandler(path)
    if suffix == ".pdf":
        return PdfHandler(path)
    raise ValueError(f"Unsupported file type: {suffix}")
